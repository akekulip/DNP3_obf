#!/usr/bin/env python3
"""Build '8. Week 8.pptx' — run with system python3 (has python-pptx)."""
from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
FIG="/home/philip/Projects/DNP3/slides_week8/figs"
RED=RGBColor(0xD0,0x00,0x00); RED2=RGBColor(0xE0,0x00,0x00)
BLUE=RGBColor(0x1F,0x4E,0x79); DARK=RGBColor(0x22,0x22,0x22); GREY=RGBColor(0x6b,0x74,0x7c)
CODEBG=RGBColor(0xF1,0xF3,0xF6); CODEFG=RGBColor(0x1b,0x2b,0x34); AMBER=RGBColor(0xB2,0x6A,0x2F)
WHITE=RGBColor(0xFF,0xFF,0xFF); BODY="Arial"; MONO="Consolas"
prs=Presentation(); prs.slide_width=In(13.333); prs.slide_height=In(7.5)
BLANK=prs.slide_layouts[6]
def slide(): return prs.slides.add_slide(BLANK)
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(In(x),In(y),In(w),In(h)); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1); return tb,tf
def setrun(r,txt,sz,color,bold=False,font=BODY,italic=False):
    r.text=txt; r.font.size=Pt(sz); r.font.color.rgb=color; r.font.bold=bold; r.font.name=font; r.font.italic=italic
def title(s,txt,color=RED,sz=34,y=0.35,x=0.6,w=12.1):
    tb,tf=box(s,x,y,w,1.0); setrun(tf.paragraphs[0].add_run(),txt,sz,color,bold=True); return tb
def subhead(s,txt,x,y,w,color=RED2,sz=15):
    tb,tf=box(s,x,y,w,0.4); setrun(tf.paragraphs[0].add_run(),txt,sz,color,bold=True)
def bullets(s,items,x,y,w,h,sz=16,gap=6,color=DARK):
    tb,tf=box(s,x,y,w,h); first=True
    for it in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(gap); p.line_spacing=1.05
        if isinstance(it,str): setrun(p.add_run(),"•  "+it,sz,color)
        else:
            setrun(p.add_run(),"•  ",sz,color)
            for t,c,b in it: setrun(p.add_run(),t,sz,c,bold=b)
    return tb
def codebox(s,lines,x,y,w,h,sz=12):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,In(x),In(y),In(w),In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=CODEBG; sh.line.color.rgb=RGBColor(0xD6,0xDD,0xE0); sh.line.width=Pt(1); sh.shadow.inherit=False
    tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Pt(8); tf.margin_right=Pt(8); tf.margin_top=Pt(6); tf.margin_bottom=Pt(6); first=True
    for ln in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False; p.line_spacing=1.06
        setrun(p.add_run(),ln,sz,CODEFG,font=MONO)
    return sh
def figure(s,name,x,y,w): return s.shapes.add_picture(os.path.join(FIG,name),In(x),In(y),width=In(w))
def footnote(s,txt,x=0.6,y=6.95,w=12.1,color=GREY,sz=11):
    tb,tf=box(s,x,y,w,0.4); setrun(tf.paragraphs[0].add_run(),txt,sz,color,italic=True)

s=slide()
tb,tf=box(s,0.9,2.3,11.5,2.6)
setrun(tf.paragraphs[0].add_run(),"Week 8",54,RED,bold=True)
setrun(tf.add_paragraph().add_run(),"14 / 07 / 2026",26,DARK)
p3=tf.add_paragraph(); p3.space_before=Pt(14)
setrun(p3.add_run(),"Binary indexes  ·  nonexistent-index experiments  ·  when to split or pad  ·  ACK delay",16,BLUE)

s=slide(); title(s,"Pick up from last week")
bullets(s,[
  [("Last week: one SELECT / OPERATE can carry ",DARK,False),("many CROBs",BLUE,True),(".  ",DARK,False),("Nmax = 16",RED2,True),(".",DARK,False)],
  "At 17 CROBs, the 17th (index 16) is rejected with TOO_MANY_OPS, so no OPERATE is sent.",
  [("Dr. Lin asked two things:",DARK,True)],
], 0.7,1.5,12.0,2.1, sz=18, gap=8)
bullets(s,[
  "Is \"too many operations\" the same as \"that index does not exist\"?",
  "Can a nonexistent index be slipped in as harmless padding?",
], 1.3,3.35,11.0,1.3, sz=17, gap=6, color=RED2)
subhead(s,"This week",0.7,4.8,4.0)
bullets(s,[
  "Binary indexes — what they are.",
  "Moved a nonexistent index around and read the errors.",
  "When do we split, and when do we pad?",
  "ACK delay — what we are thinking.",
], 0.7,5.25,12.0,1.7, sz=16, gap=4)

s=slide(); title(s,"Binary indexes")
subhead(s,"What it clarified",0.6,1.35,5.5)
bullets(s,[
  "The index is the output point a CROB targets.",
  "We simulate two binary output points: index 0 and index 1.",
  "Valid indexes run 0 to K−1. Anything K or higher does not exist.",
], 0.6,1.75,5.9,2.1, sz=15, gap=6)
subhead(s,"How I used it",0.6,3.95,5.5)
bullets(s,[
  "Configure K valid points, then ask for one extra (index K) to find the edge.",
  "Read the per-index status from the PCAP, not from the API.",
], 0.6,4.35,5.9,1.8, sz=15, gap=6)
figure(s,"fig1_binary_indexes.png",6.9,1.55,6.1)

s=slide(); title(s,"Experiments — moving the nonexistent index")
rows=[("Case","Ordered CROBs (K = 5 valid)","What we test"),
      ("End","0, 1, 2, 3, 4, 5","fake index last"),
      ("Begin","5, 0, 1, 2, 3, 4","fake index first"),
      ("Middle","0, 1, 5, 2, 3, 4","fake index in the middle"),
      ("Multiple","0, 1, 2, 3, 4, 5, 6, 7","several fake indexes"),
      ("Decoy only","5, 6, 7","all indexes fake"),
      ("Fake + limit","K=5, 17 CROBs","fake AND over the count limit"),
      ("Over limit","K=16, 17 CROBs","all valid, one too many")]
tb=s.shapes.add_table(len(rows),3,In(0.6),In(1.5),In(7.1),In(4.7)).table
tb.columns[0].width=In(1.9); tb.columns[1].width=In(3.1); tb.columns[2].width=In(2.1)
for ci in range(3):
    c=tb.cell(0,ci); c.fill.solid(); c.fill.fore_color.rgb=RED; c.vertical_anchor=MSO_ANCHOR.MIDDLE
    setrun(c.text_frame.paragraphs[0].add_run(),rows[0][ci],12,WHITE,bold=True)
for ri in range(1,len(rows)):
    for ci in range(3):
        c=tb.cell(ri,ci); c.fill.solid(); c.fill.fore_color.rgb=WHITE if ri%2 else RGBColor(0xF2,0xF5,0xF8); c.vertical_anchor=MSO_ANCHOR.MIDDLE
        col=BLUE if ci==0 else (CODEFG if ci==1 else DARK)
        setrun(c.text_frame.paragraphs[0].add_run(),rows[ri][ci],12,col,bold=(ci==0),font=MONO if ci==1 else BODY)
codebox(s,["one command, ordered CROBs:","",
  'run_master.py --action multi-crob-sbo \\','  --crob-plan "0:LATCH_ON,1:LATCH_ON,',
  '  2:LATCH_ON,3:LATCH_ON,','  4:LATCH_ON,5:LATCH_ON"'], 8.0,1.6,4.9,2.1, sz=11)
bullets(s,["Group 12 Var 1 CROBs only.","Software points only. No physical device.",
  "Correctness = per-index status, not packet count."], 8.0,4.05,4.9,2.2, sz=13, gap=6)

s=slide(); title(s,"Two different rejections")
bullets(s,[
  [("A nonexistent index → ",DARK,False),("OUT_OF_RANGE",RED2,True),("  (status 12).",DARK,False)],
  [("Too many operations → ",DARK,False),("TOO_MANY_OPS",AMBER,True),("  (status 8).",DARK,False)],
  "Different codes, and both show up per index on the wire.",
  [("K=5, N=17 shows both at once: ",DARK,False),("0–4 ok, 5–15 OUT_OF_RANGE, 17th over the limit.",DARK,True)],
], 0.6,1.45,6.0,2.7, sz=15, gap=8)
codebox(s,["# analyze_multicrob_pcap.py","STATUS_NAMES = {0: 'SUCCESS',","   8: 'TOO_MANY_OPS',","  12: 'OUT_OF_RANGE'}"], 0.6,4.4,5.9,1.6, sz=13)
figure(s,"fig3_two_errors.png",6.75,1.5,6.4)

s=slide(); title(s,"Where we put the fake index doesn't matter")
bullets(s,[
  "Fake index at the start, middle, or end — same result.",
  [("That index comes back ",DARK,False),("OUT_OF_RANGE (12)",RED2,True),(", wherever it sits.",DARK,False)],
  "The SELECT only partly succeeds, so the master sends no OPERATE.",
  "No valid output changes. The whole control is blocked.",
], 0.6,1.55,6.0,3.2, sz=16, gap=10)
figure(s,"fig2_positions.png",6.75,1.55,6.3)

s=slide(); title(s,"Can a nonexistent index be padding?  No.")
bullets(s,[
  "The fake index is rejected and is visible on the wire (OUT_OF_RANGE).",
  "One bad index blocks the whole OPERATE.",
  "So we can't slip fake indexes into a real control to pad its size.",
], 0.7,1.65,11.9,2.2, sz=19, gap=12)
sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,In(0.7),In(4.4),In(11.9),In(1.6))
sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(0xFB,0xF0,0xE8); sh.line.color.rgb=AMBER; sh.line.width=Pt(1.5); sh.shadow.inherit=False
tf=sh.text_frame; tf.word_wrap=True; tf.margin_left=Pt(14); tf.margin_top=Pt(10)
setrun(tf.paragraphs[0].add_run(),"To be clear:  ",16,AMBER,bold=True)
setrun(tf.paragraphs[0].add_run(),"this does not prove padding is impossible. It shows one obvious way (fake indexes) does not work, and points us where to look next.",16,DARK)

s=slide(); title(s,"So — when do we split, and when do we pad?")
subhead(s,"Split",0.6,1.45,5.5,color=BLUE)
bullets(s,["Cut one reply into many small packets.","Good for big reads — hides the packet shape.",
  [("But total bytes stay the same, so it ",DARK,False),("does not hide size.",RED2,True)]], 0.6,1.85,6.1,2.1, sz=15, gap=6)
subhead(s,"Pad",0.6,4.0,5.5,color=AMBER)
bullets(s,["Make a small reply look bigger — the only thing that hides size.",
  [("No safe way inside DNP3 today ",DARK,False),("(the fake-index test is one dead end).",GREY,False)],
  "Real padding = a future encrypted-tunnel step."], 0.6,4.4,6.1,2.1, sz=15, gap=6)
figure(s,"fig4_split_size.png",7.0,1.7,6.0)
footnote(s,"For now: split the big reads, normalize their timing, and record the size that still leaks.")

s=slide(); title(s,"ACK delay — what we are thinking")
bullets(s,[
  "The reply rides on the same TCP ACK that answers the request — so ACK time = reply time.",
  [("Reply time grows with N, so it leaks how many points were touched ",DARK,False),("(R²≈0.99).",RED2,True)],
  "Idea: hold each reply to a shared deadline that does not depend on N.",
  [("Random jitter won't do — an attacker averages it out. A shared clock they can't.",DARK,True)],
  "Stay under the TCP retransmit timer. Never delay a control.",
], 0.6,1.4,6.2,3.7, sz=15, gap=8)
codebox(s,["# release no earlier than a common deadline","candidate_release = max(","     response_ready_time,","     request_time + target_delay)"], 0.6,5.4,6.2,1.5, sz=13)
figure(s,"fig5_ack_timing.png",7.0,1.6,6.0)

s=slide(); title(s,"Questions / Next steps")
bullets(s,[
  [("Does the ",DARK,False),("OUT_OF_RANGE vs TOO_MANY_OPS",BLUE,True),(" story match what you expected?",DARK,False)],
  "Next: measure the real retransmit timer on the rig, and repeat the size / timing readings with error bars.",
  "Then one clean run that splits + normalizes timing, and check nothing breaks.",
  [("Open question: ",RED2,True),("do we open a future encrypted-tunnel line to actually hide size?",DARK,False)],
], 0.7,1.8,12.0,4.0, sz=19, gap=16)

out="/home/philip/Projects/DNP3/8. Week 8.pptx"
prs.save(out); print("saved",out,"| slides:",len(prs.slides._sldIdLst))
